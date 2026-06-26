"""PptxExtractor — one block per slide → slide_sequence; pictures → ImageRef (§8)."""

from __future__ import annotations

from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from trustrag.extract.plain import open_binary
from trustrag.extract.types import (
    BlockKind,
    ExtractedBlock,
    ExtractedDoc,
    ImageRef,
    SourceType,
    StructureType,
)
from trustrag.plugins.base import ExtractorPlugin


class PptxExtractor(ExtractorPlugin):
    """Each slide is one ``slide`` block (all its text frames joined), ordered as a
    sequence; ``page`` is the 1-based slide number. Pictures become ``ImageRef``s."""

    plugin_version = "pptx/1"

    def __init__(self, document_id: str | None = None) -> None:
        self.document_id = document_id

    def extract(self, source: Any) -> ExtractedDoc:
        opened, doc_id, source_uri = open_binary(source, self.document_id)
        presentation = Presentation(opened)

        blocks: list[ExtractedBlock] = []
        images: list[ImageRef] = []

        for index, slide in enumerate(presentation.slides):
            page = index + 1
            texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    frame_text = shape.text_frame.text.strip()
                    if frame_text:
                        texts.append(frame_text)
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    images.append(ImageRef(image_id=f"slide{page}-{shape.shape_id}", page=page))
            blocks.append(
                ExtractedBlock(
                    order=index,
                    kind=BlockKind.slide,
                    text="\n".join(texts),
                    page=page,
                    level=index,
                )
            )

        return ExtractedDoc(
            document_id=doc_id,
            source_type=SourceType.pptx,
            structure_type=StructureType.slide_sequence,
            source_uri=source_uri,
            blocks=tuple(blocks),
            images=tuple(images),
        )
