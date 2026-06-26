"""Shared in-memory fixtures for the extract suite — no files, no network."""

from __future__ import annotations

import io

import pytest
from docx import Document
from pptx import Presentation
from pptx.util import Inches

# A minimal valid 1x1 PNG (red pixel) — enough for an embedded image asset.
MINIMAL_PNG: bytes = (
    b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
    b"\x08\x02\x00\x00\x00\x90wS\xde\x00\x00\x00\x0cIDATx\x9cc\xf8\xcf\xc0"
    b"\x00\x00\x03\x01\x01\x00\xc9\xfe\x92\xef\x00\x00\x00\x00IEND\xaeB`\x82"
)


@pytest.fixture
def docx_bytes() -> bytes:
    """A ``.docx`` with two heading levels, body paragraphs, and one image."""
    doc = Document()
    doc.add_heading("Title One", level=1)
    doc.add_paragraph("Body para under one.")
    doc.add_heading("Sub Section", level=2)
    doc.add_paragraph("Body under sub.")
    doc.add_picture(io.BytesIO(MINIMAL_PNG))
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


@pytest.fixture
def pptx_bytes() -> bytes:
    """A ``.pptx`` with two slides; slide one carries a picture."""
    prs = Presentation()
    blank = prs.slide_layouts[6]
    slide_one = prs.slides.add_slide(blank)
    slide_one.shapes.add_textbox(
        Inches(1), Inches(1), Inches(5), Inches(1)
    ).text_frame.text = "Slide one body"
    slide_one.shapes.add_picture(
        io.BytesIO(MINIMAL_PNG), Inches(1), Inches(3), Inches(1), Inches(1)
    )
    slide_two = prs.slides.add_slide(blank)
    slide_two.shapes.add_textbox(
        Inches(1), Inches(1), Inches(5), Inches(1)
    ).text_frame.text = "Slide two body"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
