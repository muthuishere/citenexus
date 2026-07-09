"""Real DOCX and PPTX tables reach ask() as the cited passage — end to end.

Closes docs/CONTENT-COVERAGE-2026-07-08.md row 4b-prime for DOCX/PPTX: previously
a Word/PowerPoint table's text either flattened into paragraph prose (DOCX)
or was dropped entirely (PPTX, ``has_table`` never checked).
"""

from __future__ import annotations

import io
from pathlib import Path

from docx import Document
from pptx import Presentation
from pptx.util import Inches

from citenexus import CiteNexus
from citenexus.answer.result import Decision
from citenexus.testing import FakeEmbedding, FakeLLM


def _docx_with_table() -> bytes:
    doc = Document()
    doc.add_paragraph("Unrelated narrative text, outside the table.")
    table = doc.add_table(rows=3, cols=2)
    table.rows[0].cells[0].text = "Employee"
    table.rows[0].cells[1].text = "NoticeDays"
    table.rows[1].cells[0].text = "B. Singh"
    table.rows[1].cells[1].text = "45"
    table.rows[2].cells[0].text = "C. Fernandes"
    table.rows[2].cells[1].text = "60"
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _pptx_with_table() -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(
        Inches(1), Inches(0.5), Inches(5), Inches(1)
    ).text_frame.text = "Unrelated slide narrative text."
    rows, cols = 3, 2
    graphic_frame = slide.shapes.add_table(
        rows, cols, Inches(1), Inches(2), Inches(4), Inches(2)
    )
    table = graphic_frame.table
    table.cell(0, 0).text = "Employee"
    table.cell(0, 1).text = "NoticeDays"
    table.cell(1, 0).text = "B. Singh"
    table.cell(1, 1).text = "45"
    table.cell(2, 0).text = "C. Fernandes"
    table.cell(2, 1).text = "60"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_docx_table_row_is_retrieved_and_cited_verbatim(tmp_path: Path) -> None:
    docx_path = tmp_path / "notice-periods.docx"
    docx_path.write_bytes(_docx_with_table())

    rag = CiteNexus(tmp_path / "store", embedder=FakeEmbedding(), generator=FakeLLM())
    result = rag.ingest(docx_path, document_id="notice-periods-docx")
    assert result.status == "ingested"

    answer = rag.ask("How many NoticeDays does B. Singh have?")

    assert answer.evidence.decision is Decision.answered
    assert answer.claims[0].supported
    assert "B. Singh" in answer.sources[0].passage
    assert "NoticeDays: 45" in answer.sources[0].passage


def test_pptx_table_row_is_retrieved_and_cited_verbatim(tmp_path: Path) -> None:
    pptx_path = tmp_path / "notice-periods.pptx"
    pptx_path.write_bytes(_pptx_with_table())

    rag = CiteNexus(tmp_path / "store", embedder=FakeEmbedding(), generator=FakeLLM())
    result = rag.ingest(pptx_path, document_id="notice-periods-pptx")
    assert result.status == "ingested"

    answer = rag.ask("How many NoticeDays does B. Singh have?")

    assert answer.evidence.decision is Decision.answered
    assert answer.claims[0].supported
    assert "B. Singh" in answer.sources[0].passage
    assert "NoticeDays: 45" in answer.sources[0].passage
    assert answer.sources[0].page == 1
