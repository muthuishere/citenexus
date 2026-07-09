"""Document-level metadata (title/author/created/page_count) — row 8.

Each extractor reads REAL metadata from the source format's own API
(PDF /Info dict, DOCX/PPTX core_properties, HTML <title>/<meta author>) —
not fabricated, not left as dead schema. Every field is best-effort: a
format either has no such concept (HTML has no page count) or the source
document simply didn't set it.
"""

from __future__ import annotations

import io

from citenexus.domain.partition import PartitionPath
from citenexus.evidence.builder import build_evidence_units
from citenexus.evidence.unit import DocumentMetadata
from citenexus.extract.docx import DocxExtractor
from citenexus.extract.html import HtmlExtractor
from citenexus.extract.pdf import PdfExtractor
from citenexus.extract.pptx import PptxExtractor
from tests.extract.fixtures.pdf_builder import build_pdf_with_metadata


def test_pdf_metadata_is_extracted() -> None:
    pdf_bytes = build_pdf_with_metadata(
        title="Employment Policy", author="HR Dept", created="D:20260709120000Z"
    )
    doc = PdfExtractor(document_id="policy").extract(io.BytesIO(pdf_bytes))

    assert doc.metadata == DocumentMetadata(
        title="Employment Policy",
        author="HR Dept",
        created="D:20260709120000Z",
        page_count=1,
    )


def test_docx_metadata_is_extracted() -> None:
    from docx import Document

    d = Document()
    d.core_properties.title = "Employment Policy"
    d.core_properties.author = "HR Dept"
    d.add_paragraph("Body text.")
    buf = io.BytesIO()
    d.save(buf)

    doc = DocxExtractor(document_id="policy").extract(io.BytesIO(buf.getvalue()))

    assert doc.metadata is not None
    assert doc.metadata.title == "Employment Policy"
    assert doc.metadata.author == "HR Dept"
    assert doc.metadata.created is not None
    # DOCX has no fixed pagination without rendering — honestly None, not guessed.
    assert doc.metadata.page_count is None


def test_pptx_metadata_is_extracted() -> None:
    from pptx import Presentation

    p = Presentation()
    p.core_properties.title = "Quarterly Deck"
    p.core_properties.author = "Sales"
    p.slides.add_slide(p.slide_layouts[6])
    p.slides.add_slide(p.slide_layouts[6])
    buf = io.BytesIO()
    p.save(buf)

    doc = PptxExtractor(document_id="deck").extract(io.BytesIO(buf.getvalue()))

    assert doc.metadata is not None
    assert doc.metadata.title == "Quarterly Deck"
    assert doc.metadata.author == "Sales"
    # page_count is the slide count — PPTX's natural analog to a "page".
    assert doc.metadata.page_count == 2


def test_html_metadata_is_extracted() -> None:
    html = (
        '<html><head><title>Employment Policy</title>'
        '<meta name="author" content="HR Dept"></head>'
        "<body><p>Body text.</p></body></html>"
    )
    doc = HtmlExtractor(document_id="policy").extract(html)

    assert doc.metadata is not None
    assert doc.metadata.title == "Employment Policy"
    assert doc.metadata.author == "HR Dept"
    # HTML has no standard created-date or page-count concept.
    assert doc.metadata.created is None
    assert doc.metadata.page_count is None


def test_document_metadata_is_carried_onto_every_evidence_unit() -> None:
    """Carried: build_evidence_units denormalizes doc.metadata onto every EU
    (same pattern as document_id) — not just extracted, reaches the store."""
    pdf_bytes = build_pdf_with_metadata(title="Employment Policy", author="HR Dept")
    doc = PdfExtractor(document_id="policy").extract(io.BytesIO(pdf_bytes))

    partition = PartitionPath.of(("workspace", "default"))
    units = build_evidence_units(doc, partition=partition, language="en")

    assert units
    for unit in units:
        assert unit.document_metadata is not None
        assert unit.document_metadata.title == "Employment Policy"
        assert unit.document_metadata.author == "HR Dept"
