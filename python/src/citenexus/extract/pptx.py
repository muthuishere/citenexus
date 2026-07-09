"""PptxExtractor — one block per slide → slide_sequence; pictures → ImageRef (§8)."""

from __future__ import annotations

from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from citenexus.evidence.unit import DocumentMetadata
from citenexus.extract.plain import open_binary
from citenexus.extract.types import (
    BlockKind,
    ExtractedBlock,
    ExtractedDoc,
    ImageRef,
    SourceType,
    StructureType,
)
from citenexus.plugins.base import ExtractorPlugin


def _pptx_metadata(presentation: Any) -> DocumentMetadata:
    """``core_properties`` — title/author/created; ``page_count`` is the
    slide count (PPTX's natural analog to a page)."""
    props = presentation.core_properties
    created = props.created.isoformat() if props.created is not None else None
    return DocumentMetadata(
        title=props.title or None,
        author=props.author or None,
        created=created,
        page_count=len(presentation.slides),
    )


class PptxExtractor(ExtractorPlugin):
    """Each slide is one ``slide`` block (all its text frames joined), ordered as a
    sequence; ``page`` is the 1-based slide number. Pictures become ``ImageRef``s."""

    plugin_version = "pptx/1"

    def __init__(self, document_id: str | None = None) -> None:
        self.document_id = document_id

    def extract(self, source: Any) -> ExtractedDoc:
        opened, doc_id, source_uri = open_binary(source, self.document_id)
        presentation = Presentation(opened)

        blocks: list[ExtractedBlock] = []
        images: list[ImageRef] = []
        image_bytes: dict[str, bytes] = {}
        order = 0

        for index, slide in enumerate(presentation.slides):
            page = index + 1
            texts: list[str] = []
            table_rows: list[tuple[tuple[str, ...], int, str]] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    frame_text = shape.text_frame.text.strip()
                    if frame_text:
                        texts.append(frame_text)
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_id = f"slide{page}-{shape.shape_id}"
                    images.append(ImageRef(image_id=image_id, page=page))
                    try:
                        blob = shape.image.blob
                    except Exception:
                        blob = None
                    if blob:
                        image_bytes[image_id] = blob
                if shape.has_table:
                    rows = [
                        [cell.text.strip() for cell in row.cells] for row in shape.table.rows
                    ]
                    if len(rows) < 2:
                        continue
                    header = tuple(rows[0])
                    for row_index, row in enumerate(rows[1:]):
                        rendered = ", ".join(
                            f"{col}: {val}" for col, val in zip(header, row, strict=False)
                        )
                        table_rows.append((header, row_index, rendered))
            blocks.append(
                ExtractedBlock(
                    order=order,
                    kind=BlockKind.slide,
                    text="\n".join(texts),
                    page=page,
                    level=index,
                )
            )
            order += 1
            for header, row_index, rendered in table_rows:
                blocks.append(
                    ExtractedBlock(
                        order=order,
                        kind=BlockKind.table,
                        text=rendered,
                        page=page,
                        level=row_index,
                        structure_path=header,
                    )
                )
                order += 1

        return ExtractedDoc(
            document_id=doc_id,
            source_type=SourceType.pptx,
            structure_type=StructureType.slide_sequence,
            source_uri=source_uri,
            metadata=_pptx_metadata(presentation),
            blocks=tuple(blocks),
            images=tuple(images),
            image_bytes=image_bytes,
        )
